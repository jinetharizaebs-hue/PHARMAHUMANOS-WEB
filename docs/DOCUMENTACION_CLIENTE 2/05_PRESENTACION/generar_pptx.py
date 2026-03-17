#!/usr/bin/env python3
"""
Script para generar presentación PowerPoint profesional
Sistema de Gestión Comercial e-business store(EBS)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def crear_presentacion():
    """Crea presentación PowerPoint con diseño profesional"""
    
    # Crear presentación
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Colores
    COLOR_PRINCIPAL = RGBColor(0, 102, 204)  # Azul profesional
    COLOR_SECUNDARIO = RGBColor(255, 153, 0)  # Naranja
    COLOR_TEXTO = RGBColor(51, 51, 51)  # Gris oscuro
    COLOR_BLANCO = RGBColor(255, 255, 255)
    
    def agregar_titulo_subtitulo(slide, titulo, subtitulo=""):
        """Agrega título y subtítulo a diapositiva"""
        title = slide.shapes.title
        title.text = titulo
        title.text_frame.paragraphs[0].font.size = Pt(54)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = COLOR_PRINCIPAL
        
        if subtitulo:
            subtitle = slide.placeholders[1]
            subtitle.text = subtitulo
            subtitle.text_frame.paragraphs[0].font.size = Pt(24)
            subtitle.text_frame.paragraphs[0].font.color.rgb = COLOR_TEXTO
    
    def agregar_contenido(slide, titulo, puntos):
        """Agrega contenido con viñetas"""
        # Título
        title = slide.shapes.title
        title.text = titulo
        title.text_frame.paragraphs[0].font.size = Pt(44)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = COLOR_PRINCIPAL
        
        # Contenido
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9)
        height = Inches(5.5)
        
        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        
        for i, punto in enumerate(puntos):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            p.text = punto
            p.level = 0
            p.font.size = Pt(20)
            p.font.color.rgb = COLOR_TEXTO
            p.space_before = Pt(12)
            p.space_after = Pt(12)
    
    # ========== SLIDE 1: PORTADA ==========
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide1.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_PRINCIPAL
    
    # Título principal
    title_box = slide1.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Sistema de Gestión Comercial"
    title_frame.paragraphs[0].font.size = Pt(54)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = COLOR_BLANCO
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Subtítulo
    subtitle_box = slide1.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "e-business store(EBS)"
    subtitle_frame.paragraphs[0].font.size = Pt(48)
    subtitle_frame.paragraphs[0].font.bold = True
    subtitle_frame.paragraphs[0].font.color.rgb = COLOR_SECUNDARIO
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Pie de página
    footer_box = slide1.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.8))
    footer_frame = footer_box.text_frame
    footer_frame.text = "Digitalización para San Andresitos | 2 de febrero de 2026"
    footer_frame.paragraphs[0].font.size = Pt(16)
    footer_frame.paragraphs[0].font.color.rgb = COLOR_BLANCO
    footer_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # ========== SLIDE 2: CONTEXTO DEL PROBLEMA ==========
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    agregar_contenido(slide2, "El Problema Actual", [
        "❌ Facturación manual con errores frecuentes",
        "❌ Pérdida de facturas físicas (15% mensuales)",
        "❌ Cobros no sistematizados sin seguimiento",
        "❌ Inventario desactualizado",
        "❌ Cero respaldos de información crítica"
    ])
    
    # ========== SLIDE 3: IMPACTO FINANCIERO ==========
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    agregar_contenido(slide3, "Impacto Financiero Actual", [
        "💰 Pérdidas mensuales: $6.3 MILLONES COP",
        "",
        "📉 Desglose:",
        "   • Facturas extraviadas: $2.5M/mes",
        "   • Morosidad en créditos: $3M/mes",
        "   • Errores de cálculo: $800K/mes"
    ])
    
    # ========== SLIDE 4: OBJETIVO ==========
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    agregar_contenido(slide4, "Objetivo del Proyecto", [
        "✅ Digitalizar facturación electrónica",
        "✅ Automatizar seguimiento de cobros",
        "✅ Sincronizar inventario en tiempo real",
        "✅ Generar respaldos automáticos",
        "✅ Permitir acceso remoto seguro 24/7"
    ])
    
    # ========== SLIDE 5: SOLUCIÓN EBS (MÓDULOS) ==========
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    agregar_contenido(slide5, "Solución: 6 Módulos Integrados", [
        "1️⃣ Facturación Electrónica con numeración única",
        "2️⃣ Gestión de Cuentas de Cobro",
        "3️⃣ Control de Inventario en Tiempo Real",
        "4️⃣ Catálogo Digital de Productos",
        "5️⃣ Base de Datos de Clientes",
        "6️⃣ Reportes Ejecutivos"
    ])
    
    # ========== SLIDE 6: BENEFICIOS TANGIBLES ==========
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])
    agregar_contenido(slide6, "Beneficios Cuantificables", [
        "🎯 90% menos errores en facturación",
        "🎯 40% reducción en morosidad",
        "🎯 100% de facturas respaldadas",
        "🎯 80% menos tiempo en gestión",
        "🎯 ROI: Retorno de inversión en 1.85 meses"
    ])
    
    # ========== SLIDE 7: ARQUITECTURA TÉCNICA ==========
    slide7 = prs.slides.add_slide(prs.slide_layouts[5])
    agregar_contenido(slide7, "Arquitectura Técnica", [
        "🔧 Frontend: React.js + Diseño responsivo",
        "🔧 Backend: Supabase (PostgreSQL)",
        "🔧 Hosting: Vercel (alta disponibilidad)",
        "🔧 Seguridad: SSL/TLS + Encriptación",
        "🔧 Respaldos: Backup semanal automático"
    ])
    
    # ========== SLIDE 8: CRONOGRAMA ==========
    slide8 = prs.slides.add_slide(prs.slide_layouts[5])
    agregar_contenido(slide8, "Cronograma de Implementación", [
        "📅 Fase 1 (Semanas 1-2): Análisis y Preparación",
        "📅 Fase 2 (Semanas 3-4): Implementación y Migración",
        "📅 Fase 3 (Semanas 5-6): Capacitación y Pruebas",
        "📅 Fase 4 (Semanas 7-8): Puesta en Marcha",
        "",
        "⏱️ TIEMPO TOTAL: 8 semanas"
    ])
    
    # ========== SLIDE 9: INVERSIÓN ==========
    slide9 = prs.slides.add_slide(prs.slide_layouts[5])
    agregar_contenido(slide9, "Inversión Requerida", [
        "💵 Implementación (pago único): $10.000.000 COP",
        "   → 50% a la firma, 50% al finalizar",
        "",
        "💵 Mensualidad (soporte completo): $400.000 COP",
        "   → Hosting, backup, soporte, mantenimiento",
        "",
        "📊 Punto de equilibrio: 1.85 meses"
    ])
    
    # ========== SLIDE 10: COMPARATIVA ANTES/DESPUÉS ==========
    slide10 = prs.slides.add_slide(prs.slide_layouts[5])
    agregar_contenido(slide10, "Comparativa: Antes vs Después", [
        "ANTES:",
        "   • Facturas manuales: 15-20 min/factura",
        "   • Pérdidas: $6.3M/mes",
        "",
        "DESPUÉS:",
        "   • Facturación automática: 2-3 min/factura",
        "   • Recuperación: $5.4M/mes"
    ])
    
    # ========== SLIDE 11: PRÓXIMOS PASOS ==========
    slide11 = prs.slides.add_slide(prs.slide_layouts[5])
    agregar_contenido(slide11, "Próximos Pasos", [
        "1️⃣ Revisión de documentación (30 minutos)",
        "2️⃣ Reunión de aclaración (esta semana)",
        "3️⃣ Demo del sistema en vivo",
        "4️⃣ Aprobación y firma de contrato",
        "5️⃣ Inicio de implementación"
    ])
    
    # ========== SLIDE 12: CIERRE ==========
    slide12 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide12.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_PRINCIPAL
    
    closing_box = slide12.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
    closing_frame = closing_box.text_frame
    closing_frame.text = "¿Preguntas?"
    closing_frame.paragraphs[0].font.size = Pt(60)
    closing_frame.paragraphs[0].font.bold = True
    closing_frame.paragraphs[0].font.color.rgb = COLOR_BLANCO
    closing_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    contact_box = slide12.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(2))
    contact_frame = contact_box.text_frame
    contact_frame.text = "Edwin Marín\ne-business store(EBS)\n📧 contacto@distribucionesebs.com"
    contact_frame.paragraphs[0].font.size = Pt(20)
    contact_frame.paragraphs[0].font.color.rgb = COLOR_BLANCO
    contact_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Guardar presentación
    output_path = "/Users/edwinmarin/pedido-ebs-web/DOCUMENTACION_CLIENTE/05_PRESENTACION/Presentacion_EBS.pptx"
    prs.save(output_path)
    print(f"✅ Presentación creada: {output_path}")
    return output_path

if __name__ == "__main__":
    crear_presentacion()
